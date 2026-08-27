# -*- coding: utf-8 -*-
"""
合规指南 PPT 克隆器 —— 通用模板脚本
================================================
基于已有某国合规指南 .pptx 模板，克隆其版式（配色/字体/表格布局），
把文字与表格替换为目标国数据，生成 1:1 对齐的新国家合规指南。

用法
----
  1) 提取模板结构（确认形状名、表格行列数）：
     python clone_pptx.py --dump --src 模板.pptx

  2) 在下方 SPEC 中按 (slide索引 -> 形状名 -> 文本/表格行) 填写目标国数据

  3) 生成：
     python clone_pptx.py --src 模板.pptx --out 目标国_雇佣合规指南_2026.pptx

依赖：pip install python-pptx
"""
import copy
import sys

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    sys.exit(
        "缺少依赖 python-pptx。\n"
        "请先运行：pip install python-pptx\n"
        "（若使用 WorkBuddy 内置 Python，请在该环境的 pip 安装后再运行本脚本。）"
    )


# ======================================================================
# 核心工具函数（格式保留 + 表格行列对齐）
# ======================================================================

def _copy_run_format(src_run, dst_run):
    """把源 run 的全部字形属性（字号/粗体/斜体/字体/颜色）深拷贝到目标 run。"""
    src_rPr = src_run._r.find(qn('a:rPr'))
    dst_rPr = dst_run._r.find(qn('a:rPr'))
    if src_rPr is None:
        return
    if dst_rPr is None:
        dst_run._r.insert(0, copy.deepcopy(src_rPr))
        return
    # 清空目标 rPr 子元素与属性，再覆写
    for child in list(dst_rPr):
        dst_rPr.remove(child)
    for k, v in src_rPr.attrib.items():
        dst_rPr.set(k, v)
    for child in src_rPr:
        dst_rPr.append(copy.deepcopy(child))


def set_text(shape, text):
    """替换文本框内容，按 \\n 分行；逐段保留原段落的对齐方式、层级与字形格式。"""
    tf = shape.text_frame
    lines = text.split("\n")
    # 记录原每段的对齐/层级与首个非空 run（作为格式参考）
    orig = []
    for p in tf.paragraphs:
        ref = None
        for r in p.runs:
            if r.text.strip():
                ref = r
                break
        orig.append((p.alignment, p.level, ref))
    # 段落数不足则补段落
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    paras = tf.paragraphs
    for i, line in enumerate(lines):
        p = paras[i]
        al, lvl, ref = orig[i] if i < len(orig) else (None, None, None)
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = line
        if ref is not None:
            _copy_run_format(ref, run)
        if al is not None:
            p.alignment = al
        if lvl is not None:
            p.level = lvl
    # 清掉多余的段落（保留段首空行风格时不删，这里仅清文本）
    for p in paras[len(lines):]:
        for r in list(p.runs):
            r._r.getparent().remove(r._r)


def _add_table_row(table):
    """在表格末行深拷贝追加一行（保留列宽/字形），新行文本清空。"""
    tbl = table._tbl
    trs = tbl.findall(qn('a:tr'))
    last = trs[-1]
    new = copy.deepcopy(last)
    tbl.append(new)
    new_row = table.rows[-1]
    for cell in new_row.cells:
        for p in cell.text_frame.paragraphs:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.add_run().text = ""


def set_table(table, rows, tag=""):
    """把表格内容替换为 rows；先按行数对齐模板（删多余/补不足），再逐格填字。"""
    n = len(rows)
    tbl = table._tbl
    while len(table.rows) > n:
        trs = tbl.findall(qn('a:tr'))
        tbl.remove(trs[-1])
    while len(table.rows) < n:
        _add_table_row(table)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            set_text(table.cell(i, j), str(val))


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def find_table(slide, name):
    for s in slide.shapes:
        if s.has_table and s.name == name:
            return s.table
    return None


# ======================================================================
# SPEC —— 在此填写目标国数据（示例为西班牙，替换为你的国家）
#   TEXTS[slide]  = { 形状名: 多行文本 }
#   TABLES[slide] = { 表格形状名: [ [列1,列2,...], ... ] }
# ======================================================================

SPEC = {
    "TEXTS": {
        # ---- 示例：封面 ----
        1: {
            "Title 1": "西班牙雇佣合规指南",
            "Subtitle 1": "2026版 · 中企出海 HR 必备",
        },
        # ---- 示例：一页要点文字 ----
        3: {
            "Text 3": "• 最低工资 SMI €1,221/月 × 14 = €17,094/年\n"
                      "• 标准工时 40h/周（37.5h 缩减案推进中，以 BOE 为准）\n"
                      "• 社保 雇主≈30.65% / 雇员≈6.50%",
        },
    },
    "TABLES": {
        # ---- 示例：某页两张表 ----
        6: {
            "Table 0": [
                ["项目", "西班牙标准", "备注"],
                ["SMI 月薪", "€1,221", "2026 RD 126/2026"],
                ["SMI 年额", "€17,094", "14 薪制"],
            ],
            "Table 1": [
                ["维度", "数值"],
                ["标准工时", "40h/周"],
                ["单日上限", "9h"],
            ],
        },
    },
}


# ======================================================================
# 诊断：打印模板结构，用于填写 SPEC
# ======================================================================

def dump_template(src):
    prs = Presentation(src)
    print(f"SLIDES: {len(prs.slides)}\n")
    for idx, slide in enumerate(prs.slides, 1):
        print(f"===== slide {idx} =====")
        for sh in slide.shapes:
            if sh.has_table:
                t = sh.table
                print(f"  [TABLE] {sh.name}  rows={len(t.rows)} cols={len(t.columns)}")
                for ri, row in enumerate(t.rows):
                    print("    r%d:" % ri, [c.text.strip().replace("\n", " ")[:20] for c in row.cells])
            elif sh.has_text_frame and sh.text.strip():
                print(f"  [TEXT ] {sh.name}: {sh.text.strip().replace(chr(10),' / ')[:60]}")


# ======================================================================
# 构建
# ======================================================================

def build(src, out):
    prs = Presentation(src)
    for idx, slide in enumerate(prs.slides, 1):
        if idx in SPEC.get("TEXTS", {}):
            for name, text in SPEC["TEXTS"][idx].items():
                sh = find_shape(slide, name)
                if sh and sh.has_text_frame:
                    set_text(sh, text)
        if idx in SPEC.get("TABLES", {}):
            for tname, rows in SPEC["TABLES"][idx].items():
                tb = find_table(slide, tname)
                if tb:
                    set_table(tb, rows, tag=f"slide{idx}/{tname}")
    prs.save(out)
    print(f"SAVED: {out}")


# ======================================================================
# 校验：行列对齐 + 原模板国文字残留检测
# ======================================================================

def verify(out, src, leak_keywords):
    op = Presentation(out)
    ap = Presentation(src)
    mismatch = []
    for i, (s, a) in enumerate(zip(op.slides, ap.slides), 1):
        ot = [sh for sh in s.shapes if sh.has_table]
        at = [sh for sh in a.shapes if sh.has_table]
        for ts, ta in zip(ot, at):
            if len(ts.table.rows) != len(ta.table.rows) or len(ts.table.columns) != len(ta.table.columns):
                mismatch.append((i, ts.name, len(ts.table.rows), len(ta.table.rows)))
    leaks = []
    for i, s in enumerate(op.slides, 1):
        blob = ""
        for sh in s.shapes:
            if sh.has_text_frame:
                blob += " " + sh.text
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        blob += " " + c.text
        for kw in leak_keywords:
            if kw in blob:
                leaks.append((i, kw))
    print("ROW/COL MISMATCH:", mismatch if mismatch else "NONE")
    print("TEMPLATE-COUNTRY LEAKS:", leaks if leaks else "NONE")
    print("SLIDES:", len(op.slides))


# ======================================================================
# 入口
# ======================================================================

if __name__ == "__main__":
    args = sys.argv[1:]
    src = None
    out = None
    dump = False
    for a in args:
        if a == "--dump":
            dump = True
        elif a == "--src":
            i = args.index(a)
            src = args[i + 1]
        elif a == "--out":
            i = args.index(a)
            out = args[i + 1]
    if dump and src:
        dump_template(src)
    elif src and out:
        build(src, out)
        # 示例校验关键字（按模板国替换，如奥地利可填 ABGB/AngG/ASVG...）
        verify(out, src, leak_keywords=[])
    else:
        print("用法:")
        print("  python clone_pptx.py --dump --src 模板.pptx")
        print("  python clone_pptx.py --src 模板.pptx --out 目标国.pptx")
